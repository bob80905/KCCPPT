"""
TODOs
1. Improve performance by skipping verses that don't need to be parsed using a while loop maybe, instead of a for loop
2. add shadow text
3. Get this to work with firefox

"""

from pptx import Presentation
from pptx.util import Pt
from pptx.util import Inches
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import argparse
from selenium import webdriver
from selenium.webdriver.common.keys import Keys
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.common.by import By
import os
import pdb
import tkinter as tk
from tkinter import ttk


all_bible_books = [
"Genesis",
"Exodus",
"Leviticus",
"Numbers",
"Deuteronomy",
"Joshua",
"Judges",
"Ruth",
"1 Samuel",
"2 Samuel",
"1 Kings",
"2 Kings",
"1 Chronicles",
"2 Chronicles",
"Ezra",
"Nehemiah",
"Esther",
"Job",
"Psalm",
"Proverbs",
"Ecclesiastes",
"Song of Solomon",
"Isaiah",
"Jeremiah",
"Lamentations",
"Ezekiel",
"Daniel",
"Hosea",
"Joel",
"Amos",
"Obadiah",
"Jonah",
"Micah",
"Nahum",
"Habakkuk",
"Zephaniah",
"Haggai",
"Zechariah",
"Malachi",
"Matthew",
"Mark",
"Luke",
"John",
"Acts",
"Romans",
"1 Corinthians",
"2 Corinthians",
"Galatians",
"Ephesians",
"Philippians",
"Colossians",
"1 Thessalonians",
"2 Thessalonians",
"1 Timothy",
"2 Timothy",
"Titus",
"Philemon",
"Hebrews",
"James",
"1 Peter",
"2 Peter",
"1 John",
"2 John",
"3 John",
"Jude",
"Revelation"
]

book_name_to_chap_ct_dict = {
"Genesis" :	50,
"Exodus" :	40,
"Leviticus" :	27,
"Numbers" :	36,
"Deuteronomy" :	34,
"Joshua" :	24,
"Judges" :	21,
"Ruth" :	4,
"1 Samuel":	31,
"2 Samuel" :	24,
"1 Kings" :	22,
"2 Kings" :	25,
"1 Chronicles" :	29,
"2 Chronicles" :	36,
"Ezra" :	10,
"Nehemiah" :	13,
"Esther" :	10,
"Job" :	42,
"Psalm" :	150,
"Proverbs" :	31,
"Ecclesiastes" :	12,
"Song of Solomon":	8,
"Isaiah" :	66,
"Jeremiah" :	52,
"Lamentations" :	5,
"Ezekiel" :	48,
"Daniel" :	12,
"Hosea" :	14,
"Joel" :	3,
"Amos" :	9,
"Obadiah" :	1,
"Jonah" :	4,
"Micah" :	7,
"Nahum" :	3,
"Habakkuk" :	3,
"Zephaniah" :	3,
"Haggai" :	2,
"Zechariah" :	14,
"Malachi" :	4,
"Matthew" :	28,
"Mark" :	16,
"Luke" :	24,
"John" :	21,
"Acts" :	28,
"Romans" :	16,
"1 Corinthians" :	16,
"2 Corinthians" :	13,
"Galatians" :	6,
"Ephesians" :	6,
"Philippians" :	4,
"Colossians" :	4,
"1 Thessalonians" :	5,
"2 Thessalonians" :	3,
"1 Timothy" :	6,
"2 Timothy" :	4,
"Titus" :	3,
"Philemon" :	1,
"Hebrews" :	13,
"James" :	5,
"1 Peter" :	5,
"2 Peter" :	3,
"1 John" :	5,
"2 John" :	1,
"3 John" :	1,
"Jude" :	1,
"Revelation" :	22
}

class Reference:
	def __init__(self, book, chap, verse):
		self.book = book
		self.chap = chap
		self.verse = verse
	def __str__(self):
		return "{} {}:{}".format(self.book, self.chap, self.verse)
	def len(self):
		return len(str(self))

def pixels_to_points(font_size_in_pixels, screen_dpi):
    points_per_inch = 72  # 1 inch = 72 points
    font_size_in_points = (font_size_in_pixels / screen_dpi) * points_per_inch
    return font_size_in_points

def findbody(driver):
	element = driver.find_element(By.XPATH, "html/body")
	if element:
		return element
	else:
		return False

def addVersesDataToText(allChildVerseElements, text, temp_book, temp_chap, temp_verse, ref1, ref2, final_verse_reached):
	cur_chap = temp_chap
	
	for childVerseElement in allChildVerseElements:
		wordElements = childVerseElement.find_elements(By.XPATH, "*")
		wordElementIndex = 0
		while wordElementIndex < len(wordElements):
			wordElement = wordElements[wordElementIndex]
			wordElementClassAttr = wordElement.get_attribute("class")
			to_add = ""
			if (wordElement.tag_name == "sup"):
				#skip crossreferences
				wordElementIndex += 1
				continue
			if (wordElement.tag_name == "b"):
				
				if (wordElementClassAttr == "chapter-num"):
					cur_chap = int(wordElement.text.strip())									
					temp_verse = 1
					#skip chapter numbers
					
					if (cur_chap == ref2.chap and temp_verse == ref2.verse):
						final_verse_reached = True
					to_add += "#{} ${}".format(cur_chap, temp_verse) #add verse 1, which is missing on the website
					
				elif (wordElementClassAttr.startswith("verse-num")):
					# encode verse numbers
					temp_verse = int(wordElement.text.strip())
					if (cur_chap == ref2.chap and temp_verse == ref2.verse):
						final_verse_reached = True
					to_add += "${}".format(temp_verse)
			
			# Quotations are their own word element, need to be re-analyzed
			# ex: John 3:3
			elif (wordElement.tag_name == "span" and wordElementClassAttr.startswith("woc")):
				
				newChildElements = wordElement.find_elements(By.XPATH, "*")
				tempIndex = wordElementIndex+1
				for newElt in newChildElements:
					wordElements.insert(tempIndex, newElt)
					tempIndex += 1
				wordElementIndex += 1
				continue
				
			elif (wordElement.text == ""):
				# this is meant to be a space
				to_add += " "
			
			else:
				to_add += wordElement.text

			# only add if the verse is after the start
			if (all_bible_books.index(temp_book) < all_bible_books.index(ref1.book) or \
				(all_bible_books.index(temp_book) == all_bible_books.index(ref1.book) and temp_chap < ref1.chap) or \
				(all_bible_books.index(temp_book) == all_bible_books.index(ref1.book) and temp_chap == ref1.chap and int(temp_verse) < ref1.verse)):
				wordElementIndex += 1
				continue

			text += to_add
			wordElementIndex += 1
		

		if (final_verse_reached):
			break
	
	return text, final_verse_reached, temp_verse

# TODO: make this function track the current_verse, and pass it to addVersesDataToText, and use
# this new var there instead of cur_verse_num in there.
def fetch_and_encode_text(ref1, ref2):
	url = "https://www.esv.org/{}+{}/".format(ref1.book, ref1.chap)
	driver = webdriver.Chrome()
	driver.get(url)

	body = WebDriverWait(driver, 1).until(findbody)
	body.click()
	
	text = ""
	
	temp_chap = ref1.chap
	temp_book = ref1.book
	temp_verse = ref1.verse
	final_verse_reached = False
	while (True):
		section_obj = 0
		chapterobjs = driver.find_elements(By.XPATH, "/html/body/div[1]/main/article/div[1]/*")
		for ch in chapterobjs:
			if ch.get_attribute("data-reference") == "{} {}".format(temp_book, temp_chap):
				section_obj = ch
				break
			# For 1 chapter books, data-reference should just be the book
			if (temp_book in ["Obadiah", "Philemon", "2 John", "3 John", "Jude"]):
				if ch.get_attribute("data-reference") == "{}".format(temp_book):
					section_obj = ch
					break

		if (section_obj == 0):
			url = "https://www.esv.org/{}+{}/".format(temp_book, temp_chap)
			driver.get(url)
			body = WebDriverWait(driver, 1).until(findbody)
			body.click()
			continue
			
		allChildElements = section_obj.find_elements(By.XPATH, "*");
		element_index = 0
		while (element_index < len(allChildElements)):
			element = allChildElements[element_index]
			# skip chapter titles
			if (element.tag_name == "h2"):
				element_index += 1
				continue
			# skip chapter descriptors and verse sections
			if (element.tag_name == "h3"):
				element_index += 1
				continue
			# skip images
			if (element.tag_name == "img"):
				element_index += 1
				continue
			if (element.tag_name == "p"):				
				allChildVerseElements = element.find_elements(By.XPATH, "*");
				text, final_verse_reached, temp_verse = addVersesDataToText(allChildVerseElements, text, temp_book, temp_chap, temp_verse, ref1, ref2, final_verse_reached)
			if (element.tag_name == "section" and element.get_attribute("class") == "line-group"):
				allChildLineElements = element.find_elements(By.XPATH, "*");
				for line_element in allChildLineElements:
					verse_elements = line_element.find_elements(By.XPATH, "*");
					class_attr = line_element.get_attribute("class")
					
					if (class_attr == "line"):
						newtext, final_verse_reached, temp_verse = addVersesDataToText(verse_elements, text, temp_book, temp_chap, temp_verse, ref1, ref2, final_verse_reached)
						if (newtext != text):
							text = newtext + "\n"
						
					elif (class_attr.startswith("line") and class_attr.endswith("indent")):
						newtext, final_verse_reached, temp_verse = addVersesDataToText(verse_elements, text, temp_book, temp_chap, temp_verse, ref1, ref2, final_verse_reached)
						if (newtext != text):

							text += "\t" + newtext[len(text):] + "\n"

			# check that the current book, chap, and verse match the endpoint
			if temp_book == ref2.book and temp_chap == ref2.chap and final_verse_reached:
				break
			element_index += 1

		# get the next section obj which represents the next chapter
		if (int(temp_chap) + 1 <= book_name_to_chap_ct_dict[temp_book]):
			temp_chap += 1
			temp_verse = 1
			for ch in chapterobjs:
				if ch.get_attribute("data-reference") == "{} {}".format(temp_book, temp_chap):
					section_obj = ch
					break
		else:
			if (temp_book == "Revelation"):
				print("cannot go past the end of the Bible.")
			temp_book = all_bible_books[all_bible_books.index(temp_book)+1]
			temp_chap = 1
			for ch in chapterobjs:
				if ch.get_attribute("data-reference") == "{} {}".format(temp_book, temp_chap):
					section_obj = ch
					break
		if (final_verse_reached):
			break

	driver.quit()
	return text

def add_background_pic(slide, max_length, max_width):
	left = top = Inches(0)
	pic = slide.shapes.add_picture(BACKG_IMG_PATH, left, top, Pt(max_length), Pt(max_width))

def get_title_text_from_range(ref1, ref2):
	if (ref1.book == ref2.book):
		if (ref1.chap == ref2.chap):
			return "{} {}:{} - {}".format(ref1.book, ref1.chap, ref1.verse, ref2.verse)		
		else:
			return "{} {}:{} - {}:{}".format(ref1.book, ref1.chap, ref1.verse, ref2.chap, ref2.verse)		
	else:
		return "{} {}:{} - {} {}:{}".format(ref1.book, ref1.chap, ref1.verse, ref2.book, ref2.chap, ref2.verse)		
	

def add_title_text(slide, title_font_size, max_length, ref1, ref2):
	left = Inches(0)
	width = height = Pt(title_font_size)
	txBox = slide.shapes.add_textbox(left, Pt(title_font_size), Pt(max_length), Pt(title_font_size))
	txBoxShadow = txBox.shadow
	#TODO: ADD SHADOW TEXT

	tf = txBox.text_frame
	p = tf.paragraphs[0]
	p.text = get_title_text_from_range(ref1, ref2)
	p.font.bold = True
	p.font.size = Pt(36)
	p.alignment = PP_ALIGN.CENTER

	run = p.runs[0]
	run.font.shadow = True	

def add_decorative_bar(slide, slide_number, title_font_size, body_font_size, max_length, ref1, ref2):
	title_len = len(get_title_text_from_range(ref1, ref2))
	line_length = body_font_size//2 * (title_len) # include space between both references, and the hyphen
	left_endpoint = Pt((max_length/2) - (line_length//2))
	right_endpoint = Pt((max_length/2) + (line_length//2))
	multiplier = 1
	if (slide_number == 0):
		multiplier = 2.5
	top = Pt(int(multiplier*title_font_size))
	line1=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left_endpoint, top, right_endpoint, top)
	line1.line.fill.background()
	line1.line.fill.solid()
	line1.line.fill.fore_color.rgb = RGBColor(0,0,0)
	line1.line.width = Pt(1)

def find_last_textbox(slide):
    last_textbox = None

    for shape in reversed(slide.shapes):
        if shape.has_text_frame:
            last_textbox = shape
            break

    return last_textbox

def add_final_decorative_bar(slide, slide_number, title_font_size, body_font_size, max_length, max_width, ref1, ref2):
	last_textbox = find_last_textbox(slide)
	num_lines = len(last_textbox.text.split("\n"))
	if (slide_number == 0):
		space_to_textbox = int(4*title_font_size)
	else:
		space_to_textbox = int(2*title_font_size)

	print(num_lines)
	title_len = len(get_title_text_from_range(ref1, ref2))
	line_length = body_font_size//2 * (title_len) # include space between both references, and the hyphen
	left_endpoint = Pt((max_length/2) - (line_length//2))
	right_endpoint = Pt((max_length/2) + (line_length//2))
	multiplier = 1
	if (slide_number == 0):
		multiplier = 2.5
	top = Pt(int(space_to_textbox + ((num_lines + 3) * (body_font_size))))
	line1=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left_endpoint, top, right_endpoint, top)
	line1.line.fill.background()
	line1.line.fill.solid()
	line1.line.fill.fore_color.rgb = RGBColor(0,0,0)
	line1.line.width = Pt(1)

# make a list of words, where between each element of the list we can expect
# there to be a line break or space
def split_words(text):
	ret = []
	char_index = 0
	word_start = 0
	while (char_index < len(text)):
		char = text[char_index]
		if char == " ":
			ret.append(text[word_start:char_index])
			word_start = char_index+1
		if char in ["$", "#"]:
			if (word_start == char_index):
				char_index += 1
				continue

			ret.append(text[word_start:char_index])	
			word_start = char_index
			continue

		char_index += 1
	# must add the last word, since it likely isn't followed by a space
	ret.append(text[word_start:])
	return ret


def get_line(text_box_length, body_font_size, text):
	# This is a magic number assuming body font size of 36 AND slide length of 1280
	# We want 75 characters max on a single line
	remaining_width = 2420
	line = ""
	# do not actually change the line to add tabs / superscripts, just simulate how much space
	# those changes would take

	word_index = 0
	newline_found = False
	# figure out why Genesis 1 1-16, the last slide has lines that end so far away from the text box
	#make sure to pdb set trace it 
	pdb.set_trace()
	words = split_words(text)
	
	while (word_index < len(words)):
		word = words[word_index]
		char_index = 0

		while (char_index < len(word)):
			char = word[char_index]
			if (char == "\t"):
				remaining_width -= 44 # approximately the size of an indent
			elif (char == "$"):
				char_index += 1
				remaining_width -= 18 # the size of one superscripted character
				continue
			elif (char == "#"):
				char_index += 1
				remaining_width -= 38 # the size of one bolded character
				continue
			elif (char == "." or char == "," or char == "i" or char == "I" or char == "J" or char == "f" or char == "l" or char == "j" or char == "t"):
				char_index += 1
				remaining_width -= 18 # the size of one tiny character
				continue
			elif (char == "?" or char == "\""):
				char_index += 1
				remaining_width -= 24 # the size of one small character
				continue
			elif (char == "\n"):
				newline_found = True
				break
			elif (char.isupper() or char == "-" or char == "m" or char == "w"):
				char_index += 1
				remaining_width -= 42 # the size of one uppercase or big character
				continue
			else:
				remaining_width -= body_font_size 

			if (remaining_width < 0):
				break

			char_index += 1

		# firstly, if a newline was found, cut the line there
		# and return the line now
		if (newline_found):
			
			line += word[:char_index]
			rest_of_word = []
			if (len(word) > char_index):
				rest_of_word = [word[char_index+1:]]

			rest_of_text = []
			if (len(words) > word_index):
				rest_of_text = words[word_index+1:]

			return line, " ".join(rest_of_word + rest_of_text)

		# now that we're done with the word, account for the space
		remaining_width -= 18

		if (remaining_width < 0):
			# we can early return here
			# first construct the remaining text
			new_words = words[word_index:]
			return line, " ".join(new_words)

		line += word + " "
		word_index += 1

	return line, ""

def find_first_numeric_chunk(s):
	index = 0
	while (index < len(s)):
		if (s[index].isnumeric()):
			index += 1
		else:
			break
	return s[:index]


def add_text_to_textbox(text_frame, slide_number, text, body_font_size, text_box_length, text_box_width):
	# Attempt to add text to the textbox while keeping it within the specified width
	# return remaining text that doesn't fit in the box	
	p = text_frame.paragraphs[0]
	p.font.size = Pt(body_font_size)
	text_frame.word_wrap = True
	
	# 1 less line at the bottom to make it look decent / give some footer space
	max_lines = 13
	
	# The title takes up space, the body should have 2 less lines
	if (slide_number == 0):
		max_lines = 11
	line_count = 0
	
	while (line_count < max_lines):	
		
		line, text = get_line(text_box_length, body_font_size, text)
		#print("Line: {}, Text: {}".format(line, text))
		p.font.size = Pt(body_font_size)
		char_index = 0
		
		while (char_index < len(line)):
			# deal with special encoded characters
			char = line[char_index]
			
			if (char == "$"):
				verse_num = find_first_numeric_chunk(line[char_index+1:])
				superscript_run = p.add_run()
				superscript_run.font.size = Pt(25)
				superscript_run.font._element.set('baseline', '80000')
				superscript_run.text = verse_num
				
				# reset back to normal runs for later stuff
				normal_run = p.add_run()
				line = line[char_index+1 + len(verse_num):] # now skip the $, and the verse number
				char_index = 0			
				continue
			if (char == "#"):
				chap_num = find_first_numeric_chunk(line[char_index+1:])
				bold_run = p.add_run()
				bold_run.font.size = Pt(38)
				bold_run.font.bold = True
				bold_run.text = chap_num
				
				# reset back to normal runs for later stuff
				normal_run = p.add_run()
				line = line[char_index+1 + len(chap_num):] # now skip the #, and the chapter number
				char_index = 0			
				continue
			else:
				if (len(p.runs) == 0):
					p.add_run()
				p.runs[-1].text += char
				char_index += 1
				
				
		# if there's no remaining text, just return
		if (text == ""):
			return text

		line_count += 1
		p = text_frame.add_paragraph()
		if (len(p.runs) == 0):
			p.add_run()
	
	return text


# This function should return any remaining text that it wasn't able to write to the slide
def write_text_to_slide(slide, slide_number, title_font_size, body_font_size, max_length, max_width, text):
	if (slide_number == 0):
		top = Pt(int(4*title_font_size))
	else:
		top = Pt(int(2*title_font_size))	
	
	text_box_length = max_length - 2*body_font_size
	text_box_width = max_width - 4*title_font_size
	txBox = slide.shapes.add_textbox(Pt(body_font_size), top, Pt(text_box_length), Pt(text_box_width))

	#TODO: ADD SHADOW TEXT
	tf = txBox.text_frame	
	remaining_text = add_text_to_textbox(tf, slide_number, text, body_font_size, text_box_length, text_box_width)
	return remaining_text
	#p.font._element.set('baseline', '30000')

def fill_slide(args, prs, slide_number, text, ref1, ref2):
	blank_slide_layout = prs.slide_layouts[6]
	slide = prs.slides.add_slide(blank_slide_layout)
	max_length = 1280
	max_width = 720
	
	# first add the picture, since the order in which we add things is from background -> foreground
	add_background_pic(slide, max_length, max_width)
	
	# now add the title text if this is the first slide
	title_font_size = 44
	if (slide_number == 0):		
		add_title_text(slide, title_font_size, max_length, ref1, ref2)
	
	# add the decorative bar that belongs underneath each title
	body_font_size = 36
	add_decorative_bar(slide, slide_number, title_font_size, body_font_size, max_length, ref1, ref2)
	
	# Write the actual text
	remaining_text = write_text_to_slide(slide, slide_number, title_font_size, body_font_size, max_length, max_width, text)

	if (remaining_text == ""):
		add_final_decorative_bar(slide, slide_number, title_font_size, body_font_size, max_length, max_width, ref1, ref2)
	return remaining_text
		

def fill_slides(args, prs, text, ref1, ref2):
	slide_number = 0
	while (True):
		text = fill_slide(args, prs, slide_number, text, ref1, ref2)
		if (text == ""):
			break
		slide_number += 1

def retrieve_input(textBox):
    inputValue=textBox.get("1.0","end-1c")
    return inputValue

def on_button_click():

	if (args.kill):
		os.system('taskkill /F /IM POWERPNT.EXE')

	# Retrieve selected values from the drop-down lists
	value1 = dropdown1.get()
	value2 = int(dropdown2.get())
	value3 = int(dropdown3.get())
	value4 = dropdown4.get()
	value5 = int(dropdown5.get())
	value6 = int(dropdown6.get())
	global PPT_SAVE_PATH
	PPT_SAVE_PATH = retrieve_input(textbox1)
	global BACKG_IMG_PATH
	BACKG_IMG_PATH = retrieve_input(textbox2)
	startRef = Reference(value1, value2, value3)
	endRef = Reference(value4, value5, value6)

	# Print or perform actions with the selected values (replace this with your logic)
	text = fetch_and_encode_text(startRef, endRef)
	print("Text:\n{}".format(text))
	fill_slides(args, prs, text, startRef, endRef)
	prs.save(PPT_SAVE_PATH)
	print("Presentation saved to: {}".format(PPT_SAVE_PATH))

def on_button_click2():
	quit()

# Path to where the powerpoint file should be saved
PPT_SAVE_PATH = ""
# Path to the background image to use
BACKG_IMG_PATH = ""

parser = argparse.ArgumentParser(description='Generate power point slides to format and partition a given range of verses in ESV')
parser.add_argument('-k', '--kill', action='store_true', \
	help='set if you want to kill all powerpoint processes before writing to the new powerpoint file.')

args = parser.parse_args()

# Create the main window
root = tk.Tk()
root.title("KCC Powerpoint Presentation generator for Scripture Readings")

RefLabel1 = tk.Label(root, text="Starting Reference:")
RefLabel1.grid(row=0, column=0, padx=10, pady=10)

RefLabel2 = tk.Label(root, text="Ending Reference:")
RefLabel2.grid(row=0, column=2, padx=10, pady=10)

# Create and configure the first drop-down list
label1 = tk.Label(root, text="Book:")
label1.grid(row=1, column=0, padx=10, pady=10)
dropdown1 = ttk.Combobox(root, values=all_bible_books)
dropdown1.grid(row=1, column=1, padx=10, pady=10)
dropdown1.set("Genesis")  # Set default value

# Create and configure the second drop-down list
label2 = tk.Label(root, text="Chapter:")
label2.grid(row=2, column=0, padx=10, pady=10)
dropdown2 = ttk.Combobox(root, values=list(range(1, 150 + 1)))
dropdown2.grid(row=2, column=1, padx=10, pady=10)
dropdown2.set("1")  # Set default value

# Create and configure the third drop-down list
label3 = tk.Label(root, text="Verse:")
label3.grid(row=3, column=0, padx=10, pady=10)
dropdown3 = ttk.Combobox(root, values=list(range(1, 176 + 1)))
dropdown3.grid(row=3, column=1, padx=10, pady=10)
dropdown3.set("1")  # Set default value

# Create and configure the fourth drop-down list
label4 = tk.Label(root, text="Book:")
label4.grid(row=1, column=2, padx=10, pady=10)
dropdown4 = ttk.Combobox(root, values=all_bible_books)
dropdown4.grid(row=1, column=3, padx=10, pady=10)
dropdown4.set("Genesis")  # Set default value

# Create and configure the fifth drop-down list
label5 = tk.Label(root, text="Chapter:")
label5.grid(row=2, column=2, padx=10, pady=10)
dropdown5 = ttk.Combobox(root, values=list(range(1, 150 + 1)))
dropdown5.grid(row=2, column=3, padx=10, pady=10)
dropdown5.set("1")  # Set default value

# Create and configure the sixth drop-down list
label6 = tk.Label(root, text="Verse:")
label6.grid(row=3, column=2, padx=10, pady=10)
dropdown6 = ttk.Combobox(root, values=list(range(1, 176 + 1)))
dropdown6.grid(row=3, column=3, padx=10, pady=10)
dropdown6.set("1")  # Set default value

# Create text box for ppt target path
label7 = tk.Label(root, text="Powerpoint save location:")
label7.grid(row=4, column=0, padx=10, pady=10)
textbox1 = tk.Text(root, height=1, width=40)
textbox1.grid(row=4, column=1, padx=10, pady=10)

# Create text box for background image location
label8 = tk.Label(root, text="Background image location:")
label8.grid(row=4, column=2, padx=10, pady=10)
textbox2 = tk.Text(root, height=1, width=40)
textbox2.grid(row=4, column=3, padx=10, pady=10)


# Create and configure the "Enter" button
enter_button = tk.Button(root, text="Enter", command=on_button_click)
enter_button.grid(row=5, column=0, columnspan=2, pady=20)

# Create and configure the "Close" button
close_button = tk.Button(root, text="Close", command=on_button_click2)
close_button.grid(row=5, column=2, columnspan=2, pady=20)

prs = Presentation()
prs.slide_width = 1280 * 12700  # 16:9 aspect ratio, width expressed in English Metric Units
prs.slide_height = 720 * 12700  # 16:9 aspect ratio

pdb.set_trace()
# Start the main loop
root.mainloop()
# C:\Users\Josh\Desktop\Projects\KCCPPT\testppt.pptx
# C:\Users\Josh\Pictures\agera_r.jpg